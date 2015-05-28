# -*- coding: utf-8 -*-
#
# GRAFIQUEIRO V0.1
#
# AUTOR: Carlos A. Crespo
# FECHA: febrero de 2015
# Módulos para cargar archivos XLSX / PPTX
# Dependencias: python-pptx, openpyxl

from pptx import *
from openpyxl import *
from pptx.chart.data import ChartData

#
# Cuerpo del programa
#

class Grafiqueiro:

    def __init__(self):
        ''' Clase que contiene modelos, planillas y syntaxis'''
        self.ppt=None
        self.xls=None
        self.txt=""
        self.ruta_salida=""

    def abrir_ppt(self,ruta):
        '''Crea el objeto PPTX y establece texto de ruta de archivo de salida'''
        self.ppt=Presentation(ruta)
        self.ruta_salida = ruta[0:-5] + "_salida.pptx"

    def abrir_xls(self,ruta):
        ''' Crea el objeto XLSX'''
        self.xls=load_workbook(ruta).get_sheet_by_name("Datos")

    def identifica(self,b,c):
        ''' Identifica el objeto por el nombre y lo devuelve
        b = N° de slide
        c = Nombre '''
        for j in range(0,len(self.ppt.slides[b-1].shapes)):
            if self.ppt.slides[b-1].shapes[j].name == c:
                return self.ppt.slides[b-1].shapes[j]


    def carga_shape(self,b,c,d):
        '''Carga shape con el valor de una celda
        self.ppt = Presentacion
        b = N° de Slide
        c = Texto (celda entre comillas y en mayúsculas)
        d = Nombre del objeto'''
        shape = self.identifica(b,d)
        if shape <> None:
            shape.text = str(self.xls[c.upper()].value)
            self.ppt.save(self.ruta_salida)

    def carga_grafico(self,b,c,d):
        '''Carga grafico con valores de un rango de celdas
        self.ppt = Presentacion
        b = N° de Slide
        c = Rango de celdas
        d = Nombre del objeto'''
        grafico = self.identifica(b,d)
        pass

    def carga_tabla(self,b,c,d):
        '''Carga grafico con valores de un rango de celdas
        self.ppt = Presentacion
        b = N° de Slide
        c = Rango de celdas
        d = Nombre del objeto'''
        tabla = self.identifica(b,d)
        j = 0
        i = 0
        for row in self.xls[c]:
            i = 0
            for cell in row:
                if len(tabla.table.rows)>(j) and len(tabla.table.columns)>(i):
                    tabla.table.cell(j,i).text_frame.text = str(cell.value)
                    tabla.table.cell(j,i).text_frame.text = tabla.table.cell(j,i).text_frame.text.replace('None', '')
                i = i + 1
            j = j + 1
        self.ppt.save(self.ruta_salida)

    def carga_caja(self,b,c,d,e):
        '''Carga cajas con series a partir de un rango de celdas
        self.ppt = Presentacion
        b = N° de Slide
        c = Celda n° 1
        d = Nombre del objeto
        e = direccion (vertical / horizontal)'''
        lista = []
        celda = self.xls.cell(c)
        for j in range(0,len(self.ppt.slides[b-1].shapes)):
            lista.append(d+" "+str(j+1))

        ' Utilizo condicional para establecer la dirección del carga_shape
        a= 0
        if e == 1:
            for i in lista:
                self.carga_shape(b,celda.offset(a).coordinate,i)
                a += 1
        if e == 0:
            for i in lista:
                self.carga_shape(b,celda.offset(0,a).coordinate,i)
                a += 1



    def correr_syntax(ruta):
        pass



#if __name__ == '__main__':
#try:
#ejecutar()
#except (KeyboardInterrupt, SystemExit):
#pass

